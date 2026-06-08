"""memex extract — turn a content file (document / image / audio / video) into text.

Routes each file to the BEST LOCALLY-AVAILABLE extractor (delegate + reuse):
  - markitdown (if installed) handles almost everything -> Markdown; preferred.
  - else per-type fallbacks with whatever you have:
      pdf            -> pdftotext (poppler)
      docx/odt/rtf/epub/html/tex -> pandoc
      pptx/ppt/odp   -> python-pptx
      images         -> tesseract (OCR)
      audio/video    -> whisper (transcription; whisper reads video via ffmpeg)
A missing tool => graceful skip with an actionable hint (never crashes).

It REFUSES anything that isn't a document / image / audio / video (executables,
archives, libraries, fonts, binaries) — those return (None, "refused: ...").

The extractor only produces RAW text; the synth `adopt` step (the LLM) then turns
that raw dump into a clean page. Everything is local — private files never leave
the machine. Stdlib only; shells out to the user's installed tools.
"""

from __future__ import annotations

import shutil
import subprocess
import tempfile
from pathlib import Path

TEXT_EXT = {".md", ".markdown", ".mdx", ".rst", ".txt", ".org", ".adoc",
            ".csv", ".tsv", ".log"}
PDF_EXT = {".pdf"}
PANDOC_EXT = {".docx", ".doc", ".odt", ".rtf", ".epub", ".html", ".htm", ".tex"}
PPTX_EXT = {".pptx", ".ppt", ".odp"}
SHEET_EXT = {".xlsx", ".xls", ".ods"}
IMAGE_EXT = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif", ".heic"}
AUDIO_EXT = {".mp3", ".wav", ".m4a", ".aac", ".flac", ".ogg", ".opus", ".wma"}
VIDEO_EXT = {".mp4", ".mov", ".mkv", ".webm", ".avi", ".m4v"}

# the allowlist — everything else (binaries) is refused
CONTENT_EXT = (TEXT_EXT | PDF_EXT | PANDOC_EXT | PPTX_EXT | SHEET_EXT
               | IMAGE_EXT | AUDIO_EXT | VIDEO_EXT)


def _have(cmd):
    return shutil.which(cmd) is not None


def _run(cmd, timeout=600):
    try:
        out = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout)
        return out.stdout if out.returncode == 0 else None
    except Exception:
        return None


def is_content(path) -> bool:
    return Path(path).suffix.lower() in CONTENT_EXT


def available() -> dict:
    """What extractors are installed right now (for doctor / status)."""
    return {t: _have(t) for t in ("markitdown", "pandoc", "pdftotext",
                                  "tesseract", "whisper", "ffmpeg")}


def _pptx_text(fp):
    # NOTE: python-pptx is a LIBRARY, so it only works if it's in memex's OWN
    # environment (the uv-tool venv) — unlike CLI tools, which work anywhere on
    # PATH. markitdown (a CLI) is the robust, venv-independent way to do pptx.
    try:
        from pptx import Presentation
        prs = Presentation(str(fp))
    except Exception:
        return None
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        for shape in slide.shapes:
            try:
                if getattr(shape, "has_text_frame", False):
                    txt = shape.text_frame.text.strip()  # robust: all paragraphs
                    if txt:
                        lines.append(txt)
                elif getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        if any(cells):
                            lines.append(" | ".join(cells))
            except Exception:
                continue
    return "\n".join(lines).strip() or None


def _whisper(fp):
    """Transcribe audio/video (whisper reads video via ffmpeg). Slow; opt-in tool."""
    d = tempfile.mkdtemp(prefix="memex-whisper-")
    _run(["whisper", str(fp), "--model", "base", "--output_format", "txt",
          "--output_dir", d], timeout=3600)
    for f in Path(d).glob("*.txt"):
        try:
            return f.read_text(errors="ignore")
        except Exception:
            return None
    return None


def extract(fp):
    """Return (text, method) on success, or (None, reason) on skip/refusal.
    `reason` is human-readable: 'refused: …' for non-content, or an install hint."""
    fp = Path(fp)
    ext = fp.suffix.lower()

    if ext in TEXT_EXT:
        try:
            return fp.read_text(errors="ignore"), "text"
        except Exception as e:
            return None, f"read error: {e}"

    if ext not in CONTENT_EXT:
        return None, "refused: not a document / image / audio / video"

    # audio/video → LOCAL transcription ONLY. Never route media through a cloud
    # uploader (markitdown's audio path can call a cloud STT API). Privacy first.
    if ext in (AUDIO_EXT | VIDEO_EXT):
        kind = "audio" if ext in AUDIO_EXT else "video"
        if _have("whisper"):
            txt = _whisper(fp)
            if txt and txt.strip():
                return txt, "whisper"
        return None, f"{kind}: install a LOCAL transcriber — `pip install openai-whisper`"

    # documents & images: markitdown (one-stop) first, else per-type fallbacks
    if _have("markitdown"):
        out = _run(["markitdown", str(fp)])
        if out and out.strip():
            return out, "markitdown"

    # per-type fallbacks with whatever is installed
    if ext in PDF_EXT:
        if _have("pdftotext"):
            out = _run(["pdftotext", "-layout", str(fp), "-"])
            if out and out.strip():
                return out, "pdftotext"
        return None, "pdf: install markitdown, or poppler (pdftotext)"

    if ext in PANDOC_EXT:
        if _have("pandoc"):
            out = _run(["pandoc", str(fp), "-t", "gfm"])
            if out and out.strip():
                return out, "pandoc"
        return None, "doc: install markitdown or pandoc"

    if ext in PPTX_EXT:
        txt = _pptx_text(fp)
        if txt:
            return txt, "python-pptx"
        return None, "pptx: install markitdown (CLI — also does xlsx/pdf/images/audio)"

    if ext in IMAGE_EXT:
        if _have("tesseract"):
            out = _run(["tesseract", str(fp), "stdout"])
            if out and out.strip():
                return out, "tesseract-ocr"
        return None, "image: install markitdown or tesseract (OCR)"

    if ext in SHEET_EXT:
        return None, "spreadsheet: install markitdown[xlsx] or `pip install openpyxl`"

    return None, "no extractor available for this type"
